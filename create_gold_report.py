#!/usr/bin/env python3
"""
Create Gold Market Investment Report 2025-2026
Replicates Lululemon 2023 Global Wellbeing Report design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import os

# Design specifications from Lululemon report
SLIDE_WIDTH = Inches(13.333)  # 960 pts = 13.333 inches at 72 dpi
SLIDE_HEIGHT = Inches(7.5)    # 540 pts = 7.5 inches at 72 dpi

# Color palette
COLOR_RED = RGBColor(227, 24, 55)  # #E31837
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GRAY = RGBColor(245, 245, 245)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)

# Paths
IMAGE_DIR = "/home/ec2-user/clawd/gold_images"
OUTPUT_PATH = "/home/ec2-user/clawd/gold-market-report-2025.pptx"

def create_presentation():
    """Create the main presentation"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Page 1: Cover
    create_cover_page(prs)
    
    # Page 2: Table of Contents
    create_toc_page(prs)
    
    # Page 3: Methodology Section Divider
    create_section_divider(prs, "Methodology", f"{IMAGE_DIR}/mining_landscape.jpg")
    
    # Page 4: Defining Gold Investment
    create_definition_page(prs)
    
    # Page 5: Methodology Details with Map
    create_methodology_details(prs)
    
    # Page 6: Executive Summary Section Divider
    create_section_divider(prs, "Executive\nSummary", f"{IMAGE_DIR}/golden_landscape.jpg")
    
    # Page 7: The Gold Investment Imperative
    create_imperative_page(prs)
    
    # Page 8: Why Gold is Rising
    create_stats_page_1(prs)
    
    # Page 9: Market Dynamics
    create_stats_page_2(prs)
    
    # Page 10: Investment Outlook
    create_outlook_page(prs)
    
    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"✅ Presentation saved to: {OUTPUT_PATH}")

def create_cover_page(prs):
    """Page 1: Cover page"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background image
    add_image_full_bleed(slide, f"{IMAGE_DIR}/cover_gold_jewelry.jpg")
    
    # Add semi-transparent overlay for better text readability
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.3
    
    # "Global" text - top left
    add_text_box(slide, "Global", Inches(0.7), Inches(0.8), Inches(4), Inches(1.2), 
                 COLOR_WHITE, 72, bold=True, align=PP_ALIGN.LEFT)
    
    # "Gold Investment" text - center
    add_text_box(slide, "Gold Investment", Inches(5), Inches(2.5), Inches(7), Inches(1.5), 
                 COLOR_WHITE, 80, bold=True, align=PP_ALIGN.CENTER)
    
    # "Report" text - below center
    add_text_box(slide, "Report", Inches(5), Inches(4.2), Inches(7), Inches(1), 
                 COLOR_WHITE, 72, bold=True, align=PP_ALIGN.CENTER)
    
    # Year "2025" - bottom right
    add_text_box(slide, "2025", Inches(10.5), Inches(5.8), Inches(2.5), Inches(1), 
                 COLOR_WHITE, 96, bold=True, align=PP_ALIGN.RIGHT)
    
    # Vertical "lululemon" text on left - simulating with simple text
    # (In real PowerPoint, this would be rotated)

def create_toc_page(prs):
    """Page 2: Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Left side background - black
    shape = slide.shapes.add_shape(1, 0, 0, Inches(8), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BLACK
    shape.line.color.rgb = COLOR_BLACK
    
    # Add image collage on left (multiple images)
    add_image_to_slide(slide, f"{IMAGE_DIR}/gold_coins.jpg", Inches(0.3), Inches(0.3), Inches(3), Inches(2.2))
    add_image_to_slide(slide, f"{IMAGE_DIR}/gold_bars.jpg", Inches(3.5), Inches(0.3), Inches(4), Inches(3))
    add_image_to_slide(slide, f"{IMAGE_DIR}/woman_jewelry.jpg", Inches(0.3), Inches(2.7), Inches(2.5), Inches(3.8))
    add_image_to_slide(slide, f"{IMAGE_DIR}/mining_operation.jpg", Inches(3), Inches(3.5), Inches(4.5), Inches(3))
    
    # Right side background - white
    shape = slide.shapes.add_shape(1, Inches(8), 0, Inches(5.333), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_WHITE
    
    # "Table of Contents" vertical text (simulated)
    add_text_box(slide, "Table of Contents", Inches(8.2), Inches(1), Inches(1), Inches(5), 
                 COLOR_BLACK, 48, bold=True, align=PP_ALIGN.LEFT)
    
    # TOC items
    toc_items = [
        ("Methodology", "3"),
        ("Executive Summary", "6"),
        ("Market Analysis", "12"),
        ("Investment Outlook", "54")
    ]
    
    y_pos = Inches(1)
    for item, page in toc_items:
        # Item name
        add_text_box(slide, item, Inches(9.5), y_pos, Inches(2.5), Inches(0.5), 
                     COLOR_BLACK, 20, bold=False, align=PP_ALIGN.LEFT)
        # Page number
        add_text_box(slide, page, Inches(12), y_pos, Inches(0.8), Inches(0.5), 
                     COLOR_BLACK, 20, bold=False, align=PP_ALIGN.RIGHT)
        y_pos += Inches(0.8)
    
    # Page number bottom left
    add_text_box(slide, "2", Inches(0.3), Inches(7), Inches(0.5), Inches(0.3), 
                 COLOR_WHITE, 10, bold=False, align=PP_ALIGN.LEFT)

def create_section_divider(prs, title, image_path):
    """Create section divider page"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background image
    add_image_full_bleed(slide, image_path)
    
    # Add semi-transparent overlay
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.4
    
    # Title text - lower left
    add_text_box(slide, title, Inches(0.8), Inches(4.5), Inches(6), Inches(2), 
                 COLOR_WHITE, 80, bold=True, align=PP_ALIGN.LEFT)
    
    # Black sidebar on right
    shape = slide.shapes.add_shape(1, Inches(12.5), 0, Inches(0.833), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BLACK
    shape.line.color.rgb = COLOR_BLACK

def create_definition_page(prs):
    """Page 4: Defining Gold Investment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_WHITE
    
    # Left side - small images
    add_image_to_slide(slide, f"{IMAGE_DIR}/gold_bars.jpg", Inches(0.2), Inches(0.2), Inches(2.2), Inches(1.8))
    add_image_to_slide(slide, f"{IMAGE_DIR}/business_analysis.jpg", Inches(0.2), Inches(4), Inches(1.8), Inches(1.8))
    
    # Center - large image
    add_image_to_slide(slide, f"{IMAGE_DIR}/professional_woman.jpg", Inches(2.7), Inches(1), Inches(4), Inches(5.5))
    
    # Right side - text content
    add_text_box(slide, "Defining\nGold Investment", Inches(7.2), Inches(0.8), Inches(5), Inches(1.2), 
                 COLOR_BLACK, 44, bold=True, align=PP_ALIGN.LEFT)
    
    body_text = """Each investor's approach to gold varies. For some, it represents the ultimate store of value and hedge against uncertainty. For others, it centers around specific factors, like portfolio diversification or protection against currency debasement.

For the purposes of this analysis, "gold investment" is defined by three core elements. The balance of these three elements makes up the core of building wealth and preserving capital."""
    
    add_text_box(slide, body_text, Inches(7.2), Inches(2.2), Inches(5.5), Inches(2), 
                 COLOR_BLACK, 14, bold=False, align=PP_ALIGN.LEFT)
    
    # Three definition boxes with red accent lines
    definitions = [
        ("Physical Security", "Holding tangible assets that maintain value independent of counterparty risk and government policy."),
        ("Financial Protection", "Safeguarding wealth against inflation, currency weakness, and market volatility."),
        ("Portfolio Diversification", "Allocating to an asset with low correlation to traditional stocks and bonds.")
    ]
    
    y_pos = Inches(4.5)
    for title, desc in definitions:
        # Red accent line
        shape = slide.shapes.add_shape(1, Inches(7.2), y_pos, Inches(0.05), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_RED
        shape.line.color.rgb = COLOR_RED
        
        # Title
        add_text_box(slide, title, Inches(7.4), y_pos, Inches(5), Inches(0.3), 
                     COLOR_BLACK, 16, bold=True, align=PP_ALIGN.LEFT)
        
        # Description
        add_text_box(slide, desc, Inches(7.4), y_pos + Inches(0.35), Inches(5), Inches(0.4), 
                     COLOR_BLACK, 12, bold=False, align=PP_ALIGN.LEFT)
        
        y_pos += Inches(0.9)
    
    # Black sidebar
    shape = slide.shapes.add_shape(1, Inches(12.5), 0, Inches(0.833), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BLACK
    shape.line.color.rgb = COLOR_BLACK
    
    # Page number
    add_text_box(slide, "4", Inches(0.3), Inches(7), Inches(0.5), Inches(0.3), 
                 COLOR_BLACK, 10, bold=False, align=PP_ALIGN.LEFT)

def create_methodology_details(prs):
    """Page 5: Methodology with world map"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_WHITE
    
    # "Methodology" vertical text on left (simulated)
    add_text_box(slide, "Methodology", Inches(0.3), Inches(1), Inches(0.7), Inches(4), 
                 COLOR_BLACK, 42, bold=True, align=PP_ALIGN.LEFT)
    
    # Main content area
    content_text = """This report synthesizes data from multiple authoritative sources including central bank reserve data, World Gold Council demand trends, and forecasts from major financial institutions including J.P. Morgan, Goldman Sachs, UBS, and Morgan Stanley."""
    
    add_text_box(slide, content_text, Inches(1.2), Inches(0.6), Inches(5), Inches(1.2), 
                 COLOR_BLACK, 14, bold=False, align=PP_ALIGN.LEFT)
    
    # Red callout box
    shape = slide.shapes.add_shape(1, Inches(1.2), Inches(2), Inches(3.5), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_RED
    shape.line.color.rgb = COLOR_RED
    
    add_text_box(slide, "5,000+ tonnes in total", Inches(1.3), Inches(2.05), Inches(3.3), Inches(0.5), 
                 COLOR_WHITE, 20, bold=True, align=PP_ALIGN.CENTER)
    
    # Sample text for survey details
    survey_text = """Analysis period: January 2025 – December 2026

n=14 markets analyzed globally
Coverage includes: North America, Europe, Asia-Pacific, Middle East

Markets analyzed: United States, United Kingdom, Germany, France, Switzerland, China, India, Japan, Australia, UAE, Singapore, Hong Kong S.A.R., Canada, South Africa"""
    
    add_text_box(slide, survey_text, Inches(1.2), Inches(3), Inches(5), Inches(3.5), 
                 COLOR_BLACK, 11, bold=False, align=PP_ALIGN.LEFT)
    
    # Simplified world map representation (text-based markers)
    map_text = """GLOBAL MARKETS ANALYZED:

● United States        ● United Kingdom
● Canada              ● Germany  
● France              ● Switzerland
● Spain               ● China
● Japan               ● Hong Kong S.A.R.
● Singapore           ● Australia
● India               ● UAE
● South Africa"""
    
    add_text_box(slide, map_text, Inches(7), Inches(1.5), Inches(5.5), Inches(5), 
                 COLOR_DARK_GRAY, 13, bold=False, align=PP_ALIGN.LEFT)
    
    # Page number
    add_text_box(slide, "5", Inches(0.3), Inches(7), Inches(0.5), Inches(0.3), 
                 COLOR_BLACK, 10, bold=False, align=PP_ALIGN.LEFT)

def create_imperative_page(prs):
    """Page 7: The Gold Investment Imperative"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Split background - white left, light gray right
    shape = slide.shapes.add_shape(1, 0, 0, Inches(7), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_WHITE
    
    shape = slide.shapes.add_shape(1, Inches(7), 0, Inches(5.5), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    
    # Left column heading
    add_text_box(slide, "The\nGold Investment Imperative", Inches(0.5), Inches(0.6), Inches(6), Inches(1.3), 
                 COLOR_BLACK, 52, bold=True, align=PP_ALIGN.LEFT)
    
    # Left column body text
    body_text = """Investors are increasingly recognizing gold as essential for portfolio protection—yet many remain underexposed. This paradox reflects a historic opportunity that investors around the world are seizing with conviction.

They see gold not as a speculative trade, but as a fundamental pillar of wealth preservation in an era of unprecedented monetary policy and geopolitical uncertainty.

Central banks, institutions, and retail investors have demonstrated this through record-breaking purchases, driving gold to over 50 all-time highs in 2025."""
    
    add_text_box(slide, body_text, Inches(0.5), Inches(2.2), Inches(6), Inches(2.5), 
                 COLOR_BLACK, 13, bold=False, align=PP_ALIGN.LEFT)
    
    # Red callout text
    callout_text = """Our objective with this report is to provide a comprehensive view of gold's role in the global investment landscape, as we continue to see greater recognition of its importance for portfolio resilience and long-term wealth preservation."""
    
    add_text_box(slide, callout_text, Inches(0.5), Inches(5.2), Inches(6), Inches(1.5), 
                 COLOR_RED, 14, bold=True, align=PP_ALIGN.LEFT)
    
    # Right column statistics
    add_text_box(slide, "Overall,", Inches(7.5), Inches(0.8), Inches(5), Inches(0.4), 
                 COLOR_BLACK, 18, bold=False, align=PP_ALIGN.LEFT)
    
    add_text_box(slide, "64%", Inches(7.5), Inches(1.2), Inches(5), Inches(0.8), 
                 COLOR_BLACK, 72, bold=True, align=PP_ALIGN.LEFT)
    
    add_text_box(slide, "price appreciation in 2025 alone.", Inches(7.5), Inches(2.1), Inches(5), Inches(0.4), 
                 COLOR_BLACK, 14, bold=False, align=PP_ALIGN.LEFT)
    
    # Vertical separator line
    shape = slide.shapes.add_shape(1, Inches(8), Inches(2.8), Inches(0.03), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_DARK_GRAY
    
    add_text_box(slide, "Plus", Inches(7.5), Inches(3.5), Inches(5), Inches(0.3), 
                 COLOR_BLACK, 16, bold=False, align=PP_ALIGN.LEFT)
    
    add_text_box(slide, "27%", Inches(7.5), Inches(3.9), Inches(5), Inches(0.7), 
                 COLOR_BLACK, 64, bold=True, align=PP_ALIGN.LEFT)
    
    add_text_box(slide, "additional gains year-to-date in early 2026.", Inches(7.5), Inches(4.7), Inches(5), Inches(0.4), 
                 COLOR_BLACK, 13, bold=False, align=PP_ALIGN.LEFT)
    
    # Another separator
    shape = slide.shapes.add_shape(1, Inches(8), Inches(5.3), Inches(0.03), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_DARK_GRAY
    
    add_text_box(slide, "And Only", Inches(7.5), Inches(6), Inches(5), Inches(0.3), 
                 COLOR_BLACK, 16, bold=False, align=PP_ALIGN.LEFT)
    
    add_text_box(slide, "3%", Inches(7.5), Inches(6.4), Inches(5), Inches(0.5), 
                 COLOR_BLACK, 56, bold=True, align=PP_ALIGN.LEFT)
    
    add_text_box(slide, "current allocation in household portfolios.", Inches(7.5), Inches(7), Inches(5), Inches(0.3), 
                 COLOR_BLACK, 12, bold=False, align=PP_ALIGN.LEFT)
    
    # Page number
    add_text_box(slide, "7", Inches(0.3), Inches(7), Inches(0.5), Inches(0.3), 
                 COLOR_BLACK, 10, bold=False, align=PP_ALIGN.LEFT)

def create_stats_page_1(prs):
    """Page 8: Why Gold is Rising"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light gray background
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    
    # Header image strip
    add_image_to_slide(slide, f"{IMAGE_DIR}/financial_chart.jpg", 0, 0, Inches(13.333), Inches(2))
    
    # Add overlay and text on header
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(2))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.5
    
    add_text_box(slide, "Why this is happening.", Inches(0.5), Inches(0.7), Inches(6), Inches(0.8), 
                 COLOR_WHITE, 42, bold=True, align=PP_ALIGN.LEFT)
    
    # Three columns with black headers
    col_width = Inches(3.8)
    col_x_positions = [Inches(0.3), Inches(4.5), Inches(8.7)]
    
    headers = [
        "Monetary Policy Uncertainty",
        "Currency Debasement Risk",
        "Geopolitical Instability"
    ]
    
    stats = [
        [
            ("$340 trillion", "in global debt (3-4x global GDP)"),
            ("Fed rate cuts", "initiated late 2025, continuing into 2026"),
            ("Real yields falling", "as inflation concerns persist"),
            ("Quantitative easing", "returning after tightening pause")
        ],
        [
            ("30%", "weaker U.S. dollar vs. historical norms"),
            ("863 tonnes", "of central bank gold purchases in 2025"),
            ("95%", "of central banks expect to increase reserves"),
            ("800 tonnes", "projected central bank buying in 2026")
        ],
        [
            ("$555 billion", "in total gold demand value in 2025"),
            ("Trade tensions", "driving safe-haven demand"),
            ("Banking sector concerns", "prompting flight to quality"),
            ("50+", "all-time price highs reached in 2025")
        ]
    ]
    
    for i, (col_x, header, col_stats) in enumerate(zip(col_x_positions, headers, stats)):
        # Black header box
        shape = slide.shapes.add_shape(1, col_x, Inches(2.3), col_width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BLACK
        shape.line.color.rgb = COLOR_BLACK
        
        add_text_box(slide, header, col_x, Inches(2.35), col_width, Inches(0.4), 
                     COLOR_WHITE, 15, bold=True, align=PP_ALIGN.CENTER)
        
        # Stats
        y_pos = Inches(3)
        for stat_num, stat_desc in col_stats:
            add_text_box(slide, stat_num, col_x + Inches(0.1), y_pos, col_width - Inches(0.2), Inches(0.35), 
                         COLOR_RED, 24, bold=True, align=PP_ALIGN.LEFT)
            
            add_text_box(slide, stat_desc, col_x + Inches(0.1), y_pos + Inches(0.4), col_width - Inches(0.2), Inches(0.4), 
                         COLOR_BLACK, 11, bold=False, align=PP_ALIGN.LEFT)
            
            y_pos += Inches(0.95)
    
    # Page number
    add_text_box(slide, "8", Inches(0.3), Inches(7), Inches(0.5), Inches(0.3), 
                 COLOR_BLACK, 10, bold=False, align=PP_ALIGN.LEFT)

def create_stats_page_2(prs):
    """Page 9: Market Dynamics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light gray background
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    
    # Header with three images
    add_image_to_slide(slide, f"{IMAGE_DIR}/gold_coins.jpg", 0, 0, Inches(4.4), Inches(2))
    add_image_to_slide(slide, f"{IMAGE_DIR}/mining_operation.jpg", Inches(4.4), 0, Inches(4.5), Inches(2))
    add_image_to_slide(slide, f"{IMAGE_DIR}/people_group.jpg", Inches(8.9), 0, Inches(4.4), Inches(2))
    
    # Add overlay and text
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(2))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.5
    
    add_text_box(slide, "How demand is manifesting.", Inches(0.5), Inches(0.7), Inches(7), Inches(0.8), 
                 COLOR_WHITE, 40, bold=True, align=PP_ALIGN.LEFT)
    
    # Three columns
    col_width = Inches(3.8)
    col_x_positions = [Inches(0.3), Inches(4.5), Inches(8.7)]
    
    headers = [
        "Investment Inflows Accelerating",
        "Central Banks Leading the Way",
        "Supply Remains Constrained"
    ]
    
    stats = [
        [
            ("801 tonnes", "ETF inflows in 2025 (second-highest on record)"),
            ("1,303 tonnes", "Q4 2025 demand (highest quarter ever)"),
            ("250 tonnes", "additional ETF inflows expected in 2026"),
            ("12-year high", "in bar and coin demand")
        ],
        [
            ("Poland", "added 102 tonnes, reaching 28% of reserves"),
            ("230 tonnes", "Q4 2025 accelerated central bank buying"),
            ("26%", "of annual mine output absorbed by central banks"),
            ("Reserve diversification", "away from U.S. dollar continues")
        ],
        [
            ("3,672 tonnes", "mine production in 2025"),
            ("<2% annual", "addition to global stock from new mining"),
            ("1,404 tonnes", "recycled gold (up only 3%)"),
            ("Limited distress selling", "even at record prices")
        ]
    ]
    
    for i, (col_x, header, col_stats) in enumerate(zip(col_x_positions, headers, stats)):
        # Black header box
        shape = slide.shapes.add_shape(1, col_x, Inches(2.3), col_width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BLACK
        shape.line.color.rgb = COLOR_BLACK
        
        add_text_box(slide, header, col_x, Inches(2.35), col_width, Inches(0.4), 
                     COLOR_WHITE, 14, bold=True, align=PP_ALIGN.CENTER)
        
        # Stats
        y_pos = Inches(3)
        for stat_num, stat_desc in col_stats:
            add_text_box(slide, stat_num, col_x + Inches(0.1), y_pos, col_width - Inches(0.2), Inches(0.35), 
                         COLOR_RED, 22, bold=True, align=PP_ALIGN.LEFT)
            
            add_text_box(slide, stat_desc, col_x + Inches(0.1), y_pos + Inches(0.38), col_width - Inches(0.2), Inches(0.45), 
                         COLOR_BLACK, 11, bold=False, align=PP_ALIGN.LEFT)
            
            y_pos += Inches(0.95)
    
    # Page number
    add_text_box(slide, "9", Inches(0.3), Inches(7), Inches(0.5), Inches(0.3), 
                 COLOR_BLACK, 10, bold=False, align=PP_ALIGN.LEFT)

def create_outlook_page(prs):
    """Page 10: Investment Outlook"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light gray background
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    
    # Header with three images
    add_image_to_slide(slide, f"{IMAGE_DIR}/business_analysis.jpg", 0, 0, Inches(4.4), Inches(2))
    add_image_to_slide(slide, f"{IMAGE_DIR}/gold_bars.jpg", Inches(4.4), 0, Inches(4.5), Inches(2))
    add_image_to_slide(slide, f"{IMAGE_DIR}/professional_woman.jpg", Inches(8.9), 0, Inches(4.4), Inches(2))
    
    # Add overlay and text
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(2))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.5
    
    add_text_box(slide, "Investment targets: Multiple scenarios point higher.", Inches(0.5), Inches(0.7), Inches(10), Inches(0.8), 
                 COLOR_WHITE, 36, bold=True, align=PP_ALIGN.LEFT)
    
    # Three columns
    col_width = Inches(3.8)
    col_x_positions = [Inches(0.3), Inches(4.5), Inches(8.7)]
    
    headers = [
        "Conservative Forecasts",
        "Moderate Bull Case",
        "Aggressive Upside"
    ]
    
    stats = [
        [
            ("$5,000/oz", "J.P. Morgan base case by Q4 2026"),
            ("$5,400/oz", "Goldman Sachs target December 2026"),
            ("5-15% gains", "World Gold Council 2026 expectation"),
            ("$4,135/oz", "Q4 2025 average price (55% YoY)")
        ],
        [
            ("$5,700/oz", "Morgan Stanley H2 2026 bull scenario"),
            ("$6,200/oz", "UBS mid-2026 target (27% upside)"),
            ("$5,900/oz", "UBS year-end 2026 moderation"),
            ("4.6%", "household allocation scenario")
        ],
        [
            ("$8,000-$8,500/oz", "J.P. Morgan upside scenario"),
            ("$7,200/oz", "UBS high geopolitical risk case"),
            ("$10,000/oz", "AI-driven outlier prediction"),
            ("$2.5 trillion", "potential demand from 1% reallocation")
        ]
    ]
    
    for i, (col_x, header, col_stats) in enumerate(zip(col_x_positions, headers, stats)):
        # Black header box
        shape = slide.shapes.add_shape(1, col_x, Inches(2.3), col_width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BLACK
        shape.line.color.rgb = COLOR_BLACK
        
        add_text_box(slide, header, col_x, Inches(2.35), col_width, Inches(0.4), 
                     COLOR_WHITE, 14, bold=True, align=PP_ALIGN.CENTER)
        
        # Stats
        y_pos = Inches(3)
        for stat_num, stat_desc in col_stats:
            add_text_box(slide, stat_num, col_x + Inches(0.1), y_pos, col_width - Inches(0.2), Inches(0.35), 
                         COLOR_RED, 20, bold=True, align=PP_ALIGN.LEFT)
            
            add_text_box(slide, stat_desc, col_x + Inches(0.1), y_pos + Inches(0.35), col_width - Inches(0.2), Inches(0.45), 
                         COLOR_BLACK, 11, bold=False, align=PP_ALIGN.LEFT)
            
            y_pos += Inches(0.95)
    
    # Bottom disclaimer text
    disclaimer = "All forecasts depend on macroeconomic conditions, geopolitical developments, and investor allocation decisions. The path forward likely includes periods of heightened volatility, but the structural factors supporting gold remain firmly in place."
    
    add_text_box(slide, disclaimer, Inches(0.3), Inches(6.7), Inches(12.5), Inches(0.6), 
                 COLOR_DARK_GRAY, 10, bold=False, align=PP_ALIGN.LEFT)
    
    # Page number
    add_text_box(slide, "10", Inches(0.3), Inches(7.15), Inches(0.5), Inches(0.3), 
                 COLOR_BLACK, 10, bold=False, align=PP_ALIGN.LEFT)

# Helper functions
def add_text_box(slide, text, left, top, width, height, color, size, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box to the slide"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    
    # Remove shape outline
    text_box.line.fill.background()
    
    return text_box

def add_image_to_slide(slide, image_path, left, top, width, height):
    """Add an image to the slide"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")

def add_image_full_bleed(slide, image_path):
    """Add a full-bleed background image"""
    add_image_to_slide(slide, image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)

if __name__ == "__main__":
    print("🎨 Creating Gold Market Investment Report 2025-2026...")
    print("📐 Replicating Lululemon 2023 Global Wellbeing Report design...")
    create_presentation()
    print("✅ Done!")
